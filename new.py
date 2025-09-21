import os
import re
from dotenv import load_dotenv
load_dotenv()
from pptx import Presentation
from pptx.util import Pt  # Import Pt for setting font sizes
import google.generativeai as genai

# --- CONFIGURATION ---
api_key = os.environ.get('GEMINI_API_KEY')  # Your API key
OUTPUT_FOLDER = "static"

def configure_ai():

    genai.configure(api_key=api_key)
    return genai.GenerativeModel("gemini-1.5-flash-latest")

# --- FILE GENERATION (PPT ONLY) ---

def save_as_ppt(markdown_text, full_path):
    """
    Generates a PowerPoint presentation with specific formatting.
    - Title Font Size: 36pt
    - Content Font Size: 24pt
    """
    try:
        prs = Presentation()
        # Split the generated content into sections for each slide using "SLIDE:" as a delimiter
        slides_content = markdown_text.strip().split('SLIDE:')

        for content in slides_content:
            if not content.strip():
                continue
            
            # The first line of each section is the title, the rest is the body
            parts = content.strip().split('\n', 1)
            # Clean asterisks from the title and body
            title_text = parts[0].strip().replace('*', '')
            body_text = parts[1].strip().replace('*', '') if len(parts) > 1 else ""

            # Use a standard Title and Content slide layout
            slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)

            # Set title text and apply formatting
            title_shape = slide.shapes.title
            title_shape.text = title_text
            title_shape.text_frame.paragraphs[0].font.size = Pt(36)

            # Set body text and apply formatting
            if slide.placeholders[1]:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = body_text
                # Set font size for all paragraphs in the body text frame
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(24)
            
        prs.save(full_path)
        print(f"✅ PPT saved with custom formatting: {full_path}")
        # Return just the filename, which the Flask app will use
        return os.path.basename(full_path)
    except Exception as e:
        print(f"❌ Error saving PPT: {e}")
        return None

# --- PARSING AND PROMPTS ---

def parse_qna_text(raw_text):
    """
    Parses the AI's Q&A response into a structured dictionary for easy display.
    """
    sections = {}
    # Use regex to split the text by the specific headings we requested in the prompt
    pattern = r"(CEO ANALYSIS|CTO ANALYSIS|CFO ANALYSIS|MARKET ANALYSIS|MODERATOR'S FINAL VERDICT)"
    parts = re.split(pattern, raw_text)
    
    if len(parts) > 1:
        # Loop through the parts, pairing headings with their content
        for i in range(1, len(parts), 2):
            key = parts[i].replace("'", "").title()  # Clean up the key (e.g., "MODERATOR'S" -> "Moderators")
            # Clean up the content, including removing asterisks
            value = parts[i+1].strip().lstrip(':').strip().replace('*', '')
            sections[key] = value
    else:
        # If the AI didn't follow the format, show the whole text as a fallback
        sections["Full Analysis"] = raw_text.replace('*', '')
    return sections

# A highly structured prompt to ensure the AI's output is predictable and easy to parse.
QNA_PROMPT_TEMPLATE = """
You MUST follow this structure. Analyze the startup idea "{idea}" using these exact headings:
CEO ANALYSIS: (Provide analysis on market fit, user acquisition, and business model.)
CTO ANALYSIS: (Provide analysis on technical feasibility, scalability, and potential roadblocks.)
CFO ANALYSIS: (Provide analysis on monetization, profitability, and funding viability.)
MARKET ANALYSIS: (Provide plausible TAM, SAM, and SOM estimates for the next 2-3 years.)
MODERATOR'S FINAL VERDICT: (Conclude with one word: GO, PIVOT, or SCRAP, followed by a one-sentence justification.)
"""

# Prompt for the pitch deck, telling the AI to use "SLIDE:" as a separator.
PITCH_DECK_PROMPT_TEMPLATE = """
Generate content for a 5-slide pitch deck based on the analysis.
You MUST use "SLIDE:" as a separator before each slide's title.
ANALYSIS: {analysis_text}
SLIDE: Title Slide: [Catchy name and one-line pitch]
SLIDE: The Problem
SLIDE: The Solution
SLIDE: Market Size
SLIDE: Business Model & Ask
"""

# --- MAIN PROCESSING FUNCTION ---
def process_business_idea(user_idea):
    """
    This is the main function called by the Flask app. It orchestrates the entire process.
    """
    try:
        model = configure_ai()
        
        # 1. Generate the main analysis
        qna_prompt = QNA_PROMPT_TEMPLATE.format(idea=user_idea)
        qna_response = model.generate_content(qna_prompt)
        qna_text = qna_response.text

        # 2. Parse the analysis into a dictionary
        analysis_sections = parse_qna_text(qna_text)
        
        # 3. Prepare the dictionary to be returned to the Flask app
        results = {
            "analysis_sections": analysis_sections,
            "ppt_path": None,  # Will be filled with the filename later
            "error": None
        }

        # 4. Generate the pitch deck only if the idea is not scrapped
        final_verdict = analysis_sections.get("Moderator'S Final Verdict", "")
        if "SCRAP" not in final_verdict.upper():
            deck_prompt = PITCH_DECK_PROMPT_TEMPLATE.format(analysis_text=qna_text)
            deck_response = model.generate_content(deck_prompt)
            
            # Create a clean filename for the PowerPoint file
            clean_filename = re.sub(r'[\\/*?:"<>|]', "", user_idea).replace(" ", "_")[:40]
            ppt_filename = f"{clean_filename}_Pitch_Deck.pptx"
            ppt_path_full = os.path.join(OUTPUT_FOLDER, ppt_filename)
            
            # Call the save function and store the returned filename in our results
            saved_filename = save_as_ppt(deck_response.text, ppt_path_full)
            results["ppt_path"] = saved_filename
        else:
            print("❌ Idea was scrapped. No pitch deck generated.")

        return results

    except Exception as e:
        # If anything goes wrong, return the error message to be displayed on the website
        print(f"❌ An error occurred in the main processing function: {e}")
        return {"error": str(e), "analysis_sections": None, "ppt_path": None}
